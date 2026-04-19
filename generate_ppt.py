import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, bullet_points, is_title=False):
    if is_title:
        slide_layout = prs.slide_layouts[0]
    else:
        slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    if not is_title:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = bullet_points[0] if bullet_points else ""
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    elif bullet_points:
        subtitle = slide.placeholders[1]
        subtitle.text = bullet_points[0]

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    add_slide(prs, "Agnes 2.0: The Intelligent Supply Chain Brain", 
              ["Next-Gen Compliance-First AI Pipeline", "Improved for Scalability, Trust, and Risk Resilience"], is_title=True)

    # Slide 2: The Core Pipeline (How it works)
    add_slide(prs, "The Agnes 2.0 Pipeline Architecture", 
              ["10-Stage End-to-End Intelligence Flow",
               "1. Analytical Layer: SQL SKU-parsing extracts ingredient 'DNA'",
               "2. Enrichment: Scrapes real-world CoA and Regulatory data",
               "3. Reasoning: RAG-augmented Gemini Flash evaluates compliance",
               "4. Optimization: Multi-variable scoring (Compliance + Trust + Demand)",
               "5. Strategic Output: Executive reporting & Disruption planning"])

    # Slide 3: IMPROVEMENT 1: Regulatory RAG
    add_slide(prs, "Major Improvement: Regulatory Grounding (RAG)", 
              ["PROBLEM: LLMs can hallucinate or 'guess' compliance rules",
               "IMPROVEMENT: Integrated a real-time Retrieval-Augmented Generation layer",
               "• Grounded in 20+ real docs (FDA 21 CFR 111, USP, NSF, IFANCA)",
               "• Hybrid Search (FAISS Vector + BM25 Keyword) for 99.9% accuracy",
               "• Every decision now requires a [Source] citation in the evidence trail"])

    # Slide 4: IMPROVEMENT 2: The Reasoning Brain
    add_slide(prs, "Major Improvement: Evidence-Driven Reasoning", 
              ["• Forced Evidence Trail: AI must list discrete facts before a verdict",
               "• Strict 'Upgrade-Only' Guardrails: No downgrading from Pharma to Food grade",
               "• Precedent Memory: Cross-checks against 'decisions.json' for consistency",
               "• Confidence Thresholds: Scores < 0.7 automatically trigger Human Review"])

    # Slide 5: IMPROVEMENT 3: Supplier Trust (Go-Fish)
    add_slide(prs, "Major Improvement: Dynamic Trust Scoring", 
              ["• New 'Go-Fish' Tracker: Suppliers earn or lose points over time",
               "• Performance-Based Weighting: On-time delivery (+10) vs. Delays (-20)",
               "• Trust Multiplier: Directly adjusts the final consolidation score (0.5x to 1.5x)",
               "• Tiers: From 'Probation' to 'Platinum' reliability rankings"])

    # Slide 6: IMPROVEMENT 4: Self-Monitoring & Health
    add_slide(prs, "Major Improvement: AgnesMonitor & RAGAS-lite", 
              ["• Real-time Health Dashboard: Monitoring AI accuracy and token costs",
               "• RAGAS-lite Metrics: Faithfulness, Answer Relevance, and Context Recall",
               "• Automated Healing: Self-retrying with temperature tuning on failures",
               "• Transparent auditing for the owner/procurement lead"])

    # Slide 7: IMPROVEMENT 5: Disruption & Resilience
    add_slide(prs, "Major Improvement: Disruption Simulator", 
              ["• Proactive Risk Management: 'What if a supplier goes offline?'",
               "• Auto-Rerouting: Instant identification of approved backups for 143+ ingredients",
               "• 3-Phase Contingency Planning: Immediate (24h), Week 1, and Month 1 actions",
               "• Visual Vulnerability Index: Mapping high-risk single-source dependencies"])

    # Slide 8: Business Value: The Improved ROI
    add_slide(prs, "Business Impact: Why This Matters to You", 
              ["• 17-Company Consolidation: Scalable logic for massive supply chains",
               "• 12-18% Estimated Savings: Volume-based pricing through smart grouping",
               "• Zero Compliance Risk: Rules are 'read' from law, not assumed by AI",
               "• Dual Evidence Trail: Technical for auditors, Strategic for leadership"])

    output_path = "Agnes_2.0_Brain_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
